"""
Export utilities ported from the original dashboard project.

Adapts the original Plotly-based build_excel / build_pptx functions
to work with the React dashboard's chartData row format:
  chartData = [{xKey: "2024-01-01", Region1: 1.5, Region2: 2.0, ...}, ...]
"""
import io
import os
import math
from datetime import datetime
from statistics import mean

import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt, Emu
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_MARKER_STYLE,
    XL_LEGEND_POSITION,
    XL_TICK_LABEL_POSITION,
)
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.xmlchemy import OxmlElement

from pptx.oxml.ns import qn

from config.countries import REGION_TRANSLATIONS

# ── Template path ─────────────────────────────────────────────────────────────
_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "pptx_template.pptx")


# ─────────────────────────────────────────────────────────────────────────────
# PMI Heatmap helpers
# ─────────────────────────────────────────────────────────────────────────────

# Color stops matching Excel's red→yellow→green 3-color scale
# (same palette used in the slide1.xml template)
_PMI_COLOR_STOPS = [
    (44, (0xF8, 0x69, 0x6B)),
    (45, (0xF9, 0x7E, 0x6F)),
    (46, (0xFA, 0x94, 0x73)),
    (47, (0xFB, 0xAA, 0x77)),
    (48, (0xFC, 0xBF, 0x7B)),
    (49, (0xFD, 0xD5, 0x7F)),
    (50, (0xFF, 0xEB, 0x84)),
    (51, (0xE9, 0xE5, 0x83)),
    (52, (0xD3, 0xDF, 0x82)),
    (53, (0xBD, 0xD8, 0x81)),
    (54, (0xA6, 0xD2, 0x7F)),
    (55, (0x90, 0xCB, 0x7E)),
    (56, (0x7A, 0xC5, 0x7D)),
    (57, (0x63, 0xBE, 0x7B)),
]


def _pmi_value_to_rgb(value):
    """Interpolate a PMI value onto the red→yellow→green colour scale.
    Returns an (R, G, B) tuple or None for missing data."""
    if value is None:
        return None
    try:
        value = float(value)
    except (TypeError, ValueError):
        return None
    if value <= _PMI_COLOR_STOPS[0][0]:
        return _PMI_COLOR_STOPS[0][1]
    if value >= _PMI_COLOR_STOPS[-1][0]:
        return _PMI_COLOR_STOPS[-1][1]
    for i in range(len(_PMI_COLOR_STOPS) - 1):
        v1, c1 = _PMI_COLOR_STOPS[i]
        v2, c2 = _PMI_COLOR_STOPS[i + 1]
        if v1 <= value <= v2:
            t = (value - v1) / (v2 - v1)
            return (
                int(c1[0] + t * (c2[0] - c1[0])),
                int(c1[1] + t * (c2[1] - c1[1])),
                int(c1[2] + t * (c2[2] - c1[2])),
            )
    return _PMI_COLOR_STOPS[-1][1]


def _remove_cell_borders(cell):
    """Set all four borders of a table cell to noFill (invisible)."""
    tc = cell._tc
    # Find or create <a:tcPr>
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    for ln_tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        existing = tcPr.find(qn(ln_tag))
        if existing is not None:
            tcPr.remove(existing)
        ln_el = etree.SubElement(tcPr, qn(ln_tag))
        etree.SubElement(ln_el, qn('a:noFill'))


_MONTHS_DE = ['Jan', 'Feb', 'Mrz', 'Apr', 'Mai', 'Jun',
              'Jul', 'Aug', 'Sep', 'Okt', 'Nov', 'Dez']


def _fmt_pmi_col_date(date_str):
    """Format a date string as German short month + 2-digit year, e.g. 'Mrz 26'."""
    try:
        d = datetime.fromisoformat(str(date_str)[:10])
        return f"{_MONTHS_DE[d.month - 1]} {str(d.year)[-2:]}"
    except Exception:
        return str(date_str)[:7]


def _add_pmi_heatmap_table(item: dict, chart_pl, slide) -> bool:
    """
    Replace the chart placeholder with a color-coded PMI heatmap table.

    The data is expected to be already aggregated to monthly resolution
    (sortedData from the frontend, newest first, max 13 months).
    Each cell background is filled using the red→yellow→green PMI scale.
    """
    chart_data = item.get('chartData') or []
    regions    = item.get('regions') or []
    x_key      = item.get('xKey') or 'DatePoint'

    if not chart_data or not regions:
        return False

    # Data is already in newest-first order (sorted by frontend)
    date_cols = chart_data
    num_date_cols = len(date_cols)
    num_regions   = len(regions)

    # Save placeholder geometry before removing it
    left   = chart_pl.left
    top    = chart_pl.top
    width  = chart_pl.width
    height = chart_pl.height

    # Remove the chart placeholder from the slide spTree
    chart_pl._element.getparent().remove(chart_pl._element)

    # Table dimensions: header row + one row per region; label col + one col per date
    n_rows = num_regions + 1
    n_cols = num_date_cols + 1

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Disable all special first/last row/col formatting
    tbl.first_row  = False
    tbl.first_col  = False
    tbl.last_row   = False
    tbl.last_col   = False

    # Column widths: label col ~15% of total, rest split equally
    label_col_w = int(width * 0.14)
    data_col_w  = (width - label_col_w) // num_date_cols
    tbl.columns[0].width = label_col_w
    for i in range(1, n_cols):
        tbl.columns[i].width = data_col_w

    # Row heights: equal
    row_h = height // n_rows
    for i in range(n_rows):
        tbl.rows[i].height = row_h

    # ── Header row ────────────────────────────────────────────────────────────
    hdr = tbl.cell(0, 0)
    hdr.text = 'Region'
    tf = hdr.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(10)
    _remove_cell_borders(hdr)

    for ci, row in enumerate(date_cols):
        cell = tbl.cell(0, ci + 1)
        cell.text = _fmt_pmi_col_date(row.get(x_key, ''))
        tf = cell.text_frame
        for para in tf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(10)
        _remove_cell_borders(cell)

    # ── Data rows ─────────────────────────────────────────────────────────────
    for ri, region in enumerate(regions):
        # Label cell
        label_cell = tbl.cell(ri + 1, 0)
        label_cell.text = REGION_TRANSLATIONS.get(region, region)
        tf = label_cell.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(10)
        _remove_cell_borders(label_cell)

        # Value cells
        for ci, date_row in enumerate(date_cols):
            raw = date_row.get(region)
            cell = tbl.cell(ri + 1, ci + 1)
            if raw is not None:
                try:
                    fval = float(raw)
                    cell.text = str(int(round(fval)))
                    rgb = _pmi_value_to_rgb(fval)
                    if rgb:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*rgb)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                except (TypeError, ValueError):
                    cell.text = ''
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                cell.text = ''
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            tf = cell.text_frame
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(10)
            _remove_cell_borders(cell)

    return True


# ─────────────────────────────────────────────────────────────────────────────
# Helpers shared by Excel and PPTX
# ─────────────────────────────────────────────────────────────────────────────

def _chartdata_to_traces(chart_data: list, regions: list, x_key: str) -> list:
    """
    Convert recharts row-format data to Plotly-like scatter trace list.
    Skips any region whose values are all non-numeric (e.g. string metadata columns).
    Handles both plain region keys ('Germany') and metric-prefixed keys
    ('Germany_BEST_PE_RATIO'), using only the prefix as the display name.
    """
    traces = []
    for region in regions:
        x_vals = [row.get(x_key) for row in chart_data]
        y_vals = [row.get(region) for row in chart_data]
        # Only include traces with at least one numeric value
        has_numeric = any(
            isinstance(v, (int, float)) and not (isinstance(v, float) and (math.isnan(v) or math.isinf(v)))
            for v in y_vals
        )
        if not has_numeric:
            continue
        # Use only the part before the first '_' as display name (e.g. "Germany_PE" → "Germany")
        base_name = region.split("_")[0] if "_" in region else region
        display_name = REGION_TRANSLATIONS.get(base_name, base_name)
        traces.append({
            "type": "scatter",
            "name": display_name,
            "x": x_vals,
            "y": y_vals,
            "showlegend": True,
        })
    return traces


def _clean_chart_values(values: list) -> list:
    """Replace NaN/Inf/non-numeric values with None so python-pptx skips them gracefully."""
    cleaned = []
    for v in values:
        if v is None:
            cleaned.append(None)
        elif isinstance(v, (int, float)):
            cleaned.append(None if (math.isnan(v) or math.isinf(v)) else v)
        else:
            # Strings and other non-numeric types → None
            cleaned.append(None)
    return cleaned


# ─────────────────────────────────────────────────────────────────────────────
# Excel export  (mirrors original build_excel, but reads our trace format)
# ─────────────────────────────────────────────────────────────────────────────

def build_excel(items: list) -> bytes:
    """
    Build an Excel workbook from a list of chart items.

    Items are grouped by their 'group' value; same group = same sheet.
    Each chart is placed side-by-side (advancing startcol) with its title
    written to row 0.

    Item shape:
      {
        "title": str,
        "group": int,
        "chartData": list[dict],   # row objects
        "regions": list[str],
        "xKey": str,
      }
    """
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        grouped: dict[int, list] = {}
        for it in items:
            grp = it.get("group", 1)
            grouped.setdefault(grp, []).append(it)

        for gnum in sorted(grouped):
            sheet_name = f"Group{gnum}"
            startcol = 0

            for it in grouped[gnum]:
                subheading   = it.get("subheading", "")
                y_axis_label = it.get("yAxisLabel", "")

                # ── Balken chart: write a summary statistics table ─────────
                if it.get("chartType") == "Bar" and it.get("balkenData"):
                    balken_items = it["balkenData"]
                    df = pd.DataFrame([
                        {
                            "Name":    d.get("name", ""),
                            "Min":     d.get("min"),
                            "Max":     d.get("max"),
                            "Aktuell": d.get("current"),
                        }
                        for d in balken_items
                    ])
                    data_start_row = 1
                    if subheading:
                        data_start_row += 1
                    if y_axis_label:
                        data_start_row += 1
                    df.to_excel(
                        writer,
                        sheet_name=sheet_name,
                        startrow=data_start_row,
                        startcol=startcol,
                        index=False,
                    )
                    worksheet = writer.sheets[sheet_name]
                    worksheet.write(0, startcol, it["title"])
                    next_row = 1
                    if subheading:
                        worksheet.write(next_row, startcol, subheading)
                        next_row += 1
                    if y_axis_label:
                        worksheet.write(next_row, startcol, f"Einheit: {y_axis_label}")
                    startcol += df.shape[1] + 1
                    continue

                # ── Standard time-series chart ─────────────────────────────
                chart_data = it.get("chartData") or []
                regions = it.get("regions") or []
                x_key = it.get("xKey") or "DatePoint"

                if not chart_data:
                    continue

                # Build traces like the original (x + y series per region)
                traces = _chartdata_to_traces(chart_data, regions, x_key)

                # Filter out flat / all-None traces (mirrors original)
                valid_traces = []
                for tr in traces:
                    y = tr.get("y", [])
                    non_null = [v for v in y if v is not None and not (isinstance(v, float) and math.isnan(v))]
                    if len(non_null) > 1 and len(set(non_null)) == 1:
                        continue   # flat line – skip
                    valid_traces.append(tr)

                if not valid_traces:
                    continue

                # Build DataFrame: x column + one column per region
                series = []
                x = valid_traces[0].get("x", [])
                series.append(pd.Series(x, name="x"))
                for idx, tr in enumerate(valid_traces, start=1):
                    nm = tr.get("name") or f"series{idx}"
                    series.append(pd.Series(tr.get("y", []), name=nm))
                df = pd.concat(series, axis=1)

                # Write to sheet: title in row 0, optional subheading in row 1,
                # optional y-axis label in row 2, data below
                data_start_row = 1
                if subheading:
                    data_start_row += 1
                if y_axis_label:
                    data_start_row += 1
                df.to_excel(
                    writer,
                    sheet_name=sheet_name,
                    startrow=data_start_row,
                    startcol=startcol,
                    index=False,
                )
                worksheet = writer.sheets[sheet_name]
                worksheet.write(0, startcol, it["title"])
                next_row = 1
                if subheading:
                    worksheet.write(next_row, startcol, subheading)
                    next_row += 1
                if y_axis_label:
                    worksheet.write(next_row, startcol, f"Einheit: {y_axis_label}")

                startcol += df.shape[1] + 1  # advance for next chart

    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# PPTX helpers  (ported from original pptx_download.py)
# ─────────────────────────────────────────────────────────────────────────────

def _set_date_format(chart, date_format: str):
    cat_axis_xml = chart.category_axis._element
    for num_fmt in cat_axis_xml.findall(".//c:numFmt", namespaces=cat_axis_xml.nsmap):
        num_fmt.set("formatCode", date_format)
        num_fmt.set("sourceLinked", "0")


def _set_value_format(chart, number_format: str):
    val_axis_xml = chart.value_axis._element
    nsmap = val_axis_xml.nsmap
    num_fmt_elems = val_axis_xml.findall(".//c:numFmt", namespaces=nsmap)
    if num_fmt_elems:
        for num_fmt in num_fmt_elems:
            num_fmt.set("formatCode", number_format)
            num_fmt.set("sourceLinked", "0")
    else:
        ns = nsmap.get("c", "http://schemas.openxmlformats.org/drawingml/2006/chart")
        num_fmt = etree.Element(f"{{{ns}}}numFmt", formatCode=number_format, sourceLinked="0")
        scaling = val_axis_xml.find("c:scaling", namespaces=nsmap)
        if scaling is not None:
            val_axis_xml.insert(list(val_axis_xml).index(scaling) + 1, num_fmt)
        else:
            val_axis_xml.insert(0, num_fmt)


def _set_precision(span: float, target_ticks: int = 10) -> float:
    if span <= 0:
        return 1.0
    raw_step = span / float(target_ticks)
    exp = math.floor(math.log10(raw_step))
    mantissa = raw_step / (10 ** exp)
    if   mantissa <= 1: nice_m = 1
    elif mantissa <= 2: nice_m = 2
    elif mantissa <= 5: nice_m = 5
    else:               nice_m = 10
    return nice_m * (10 ** exp)


def _round_to_precision(number: float, precision: float, direction: str) -> float:
    decimal_places = len(str(precision).split(".")[1]) if "." in str(precision) else 0
    if direction == "up":
        return round(number + (precision - (number % precision)) if number % precision != 0 else number, decimal_places)
    else:
        return round(number - (number % precision), decimal_places)


def _SubElement(parent, tagname: str, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_value_axis_title(chart, title_text: str):
    """Inject a y-axis title by parsing the exact XML structure from chart1.xml.

    chart1.xml has the title between <c:majorGridlines> and <c:numFmt>:
      <c:title>
        <c:tx><c:rich>
          <a:bodyPr/><a:lstStyle/>
          <a:p>
            <a:pPr><a:defRPr sz="1200" b="0"/></a:pPr>
            <a:r><a:rPr lang="de-DE" sz="1200" b="0" dirty="0"/><a:t>Wert</a:t></a:r>
          </a:p>
        </c:rich></c:tx>
        <c:overlay val="0"/>
      </c:title>
    chart2.xml has no <c:title> element at all.
    """
    if not title_text:
        return
    try:
        val_ax_xml = chart.value_axis._element
    except Exception:
        return

    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    # Remove any existing title so we never duplicate it
    existing = val_ax_xml.find(f"{{{c_ns}}}title")
    if existing is not None:
        val_ax_xml.remove(existing)

    # Build the exact XML string matching chart1.xml and parse it with lxml.
    # Including explicit namespace declarations on the root element guarantees
    # lxml serialises with the correct c: / a: prefixes when inserted into the tree.
    safe_text = (
        title_text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )
    xml_str = (
        '<c:title'
        ' xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<c:tx><c:rich>'
        '<a:bodyPr/>'
        '<a:lstStyle/>'
        '<a:p>'
        '<a:pPr><a:defRPr sz="1200" b="0"/></a:pPr>'
        '<a:r>'
        '<a:rPr lang="de-DE" sz="1200" b="0" dirty="0"/>'
        f'<a:t>{safe_text}</a:t>'
        '</a:r>'
        '</a:p>'
        '</c:rich></c:tx>'
        '<c:overlay val="0"/>'
        '</c:title>'
    )
    title_el = etree.fromstring(xml_str)

    # Insert after c:majorGridlines (if present), else before c:numFmt, else at end
    children  = list(val_ax_xml)
    gridlines = val_ax_xml.find(f"{{{c_ns}}}majorGridlines")
    num_fmt   = val_ax_xml.find(f"{{{c_ns}}}numFmt")
    if gridlines is not None:
        idx = children.index(gridlines) + 1
    elif num_fmt is not None:
        idx = children.index(num_fmt)
    else:
        idx = len(children)
    val_ax_xml.insert(idx, title_el)


def _set_cross_between(chart, value: str = "midCat"):
    val_axis_xml = chart.value_axis._element
    cross = val_axis_xml.find(".//c:crossBetween", namespaces=val_axis_xml.nsmap)
    if cross is not None:
        cross.set("val", value)
    else:
        cross = etree.Element(
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}crossBetween",
            val=value,
        )
        val_axis_xml.append(cross)


# ─────────────────────────────────────────────────────────────────────────────
# Balken (mixed bar + scatter) chart helpers
# ─────────────────────────────────────────────────────────────────────────────

def _hex_to_pptx_hex(hex_color: str) -> str:
    """Convert '#8b5cf6' → '8B5CF6' for OOXML srgbClr val attribute."""
    c = hex_color.strip().lstrip('#')
    if len(c) == 3:
        c = ''.join(ch * 2 for ch in c)
    return c.upper()[:6]


def _xml_esc(text: str) -> str:
    """Escape XML-special characters in text content."""
    return (
        str(text)
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
    )


def _add_balken_chart(bar_items: list, item: dict, chart_pl) -> bool:
    """
    Insert a mixed stacked-bar + scatter (Balken) chart into a PPTX placeholder.

    Uses chart1.xml as the structural template: a barChart (stacked columns for
    the historical min→max range) overlaid with a scatterChart (current-value dots).

    bar_items: list of dicts with keys:
        name     – category label for the X axis (series / country name)
        spacer   – transparent base (the historical minimum value)
        range    – visible bar height (historical max – min)
        current  – current value shown as a scatter dot
        color    – hex colour string, e.g. '#8b5cf6'
    """
    if not bar_items:
        return False

    C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    names      = [_xml_esc(d.get('name', ''))            for d in bar_items]
    spacers    = [float(d.get('spacer')  or 0)            for d in bar_items]
    ranges_    = [float(d.get('range')   or 0)            for d in bar_items]
    currents   = [d.get('current')                        for d in bar_items]
    medians    = [d.get('median')                         for d in bar_items]
    hex_colors = [_hex_to_pptx_hex(d.get('color', '#4472C4')) for d in bar_items]
    n          = len(bar_items)

    # Axis IDs – match chart1.xml values so shared axes work correctly
    CAT_AX_ID = "636977375"
    VAL_AX_ID = "636974975"

    # ── XML fragment builders ────────────────────────────────────────────────

    def _num_pt(idx: int, val) -> str:
        if val is None or (isinstance(val, float) and (math.isnan(val) or math.isinf(val))):
            return ''
        return f'<c:pt idx="{idx}"><c:v>{val}</c:v></c:pt>'

    def _str_pt(idx: int, val: str) -> str:
        return f'<c:pt idx="{idx}"><c:v>{val}</c:v></c:pt>'

    # Shared category cache (sector / country names on the X axis)
    cat_pts  = ''.join(_str_pt(i, name) for i, name in enumerate(names))
    cat_ref  = (
        f'<c:strRef>'
        f'<c:f>Sheet1!$A$2:$A${n + 1}</c:f>'
        f'<c:strCache>'
        f'<c:ptCount val="{n}"/>'
        f'{cat_pts}'
        f'</c:strCache>'
        f'</c:strRef>'
    )

    def _num_ref(values: list, fmt: str = "General") -> str:
        pts = ''.join(_num_pt(i, v) for i, v in enumerate(values))
        return (
            f'<c:numRef><c:f/>'
            f'<c:numCache>'
            f'<c:formatCode>{fmt}</c:formatCode>'
            f'<c:ptCount val="{n}"/>'
            f'{pts}'
            f'</c:numCache>'
            f'</c:numRef>'
        )

    # Fixed palette: all bars one colour, all dots another, all median ticks a third
    BAR_COLOR    = '4472C4'   # blue
    DOT_COLOR    = 'ED7D31'   # orange
    MEDIAN_COLOR = 'FFC000'   # gold

    # Y-axis number format: fewer decimals for large values
    all_vals = spacers + ranges_ + [v for v in currents if v is not None]
    maxabs   = max((abs(v) for v in all_vals), default=1)
    num_fmt  = "#,##0" if maxabs > 20 else "#,##0.0"

    # ── Full <c:chart> XML (mirrors chart1.xml structure) ───────────────────
    chart_xml_head = (
        f'<c:chart'
        f' xmlns:c="{C_NS}"'
        f' xmlns:a="{A_NS}">'  

        f'<c:autoTitleDeleted val="1"/>'
        f'<c:plotArea>'
        f'<c:layout/>'

        # ── barChart: two stacked series (spacer transparent + range coloured) ──
        f'<c:barChart>'
        f'<c:barDir val="col"/>'
        f'<c:grouping val="stacked"/>'
        f'<c:varyColors val="0"/>'

        # Series 0 – transparent spacer (0 → min)
        f'<c:ser>'
        f'<c:idx val="0"/><c:order val="0"/>'
        f'<c:spPr>'
        f'<a:noFill/>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        f'<c:invertIfNegative val="0"/>'
        f'<c:cat>{cat_ref}</c:cat>'
        f'<c:val>{_num_ref(spacers)}</c:val>'
        f'</c:ser>'

        # Series 1 – coloured range (min → max), uniform colour for all bars
        f'<c:ser>'
        f'<c:idx val="1"/><c:order val="1"/>'
        f'<c:spPr>'
        f'<a:solidFill><a:srgbClr val="{BAR_COLOR}"/></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        f'<c:invertIfNegative val="0"/>'
        f'<c:cat>{cat_ref}</c:cat>'
        f'<c:val>{_num_ref(ranges_)}</c:val>'
        f'</c:ser>'

        f'<c:dLbls>'
        f'<c:showLegendKey val="0"/><c:showVal val="0"/>'
        f'<c:showCatName val="0"/><c:showSerName val="0"/>'
        f'<c:showPercent val="0"/><c:showBubbleSize val="0"/>'
        f'</c:dLbls>'
        f'<c:gapWidth val="150"/>'
        f'<c:overlap val="100"/>'
        f'<c:axId val="{CAT_AX_ID}"/>'
        f'<c:axId val="{VAL_AX_ID}"/>'
        f'</c:barChart>'

        # ── scatterChart: current-value dots with per-point coloured markers ─
        f'<c:scatterChart>'
        f'<c:scatterStyle val="lineMarker"/>'
        f'<c:varyColors val="0"/>'
        f'<c:ser>'
        f'<c:idx val="2"/><c:order val="2"/>'
        f'<c:spPr>'
        f'<a:ln w="25400" cap="rnd"><a:noFill/><a:round/></a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        # Single uniform marker for all current-value dots
        f'<c:marker>'
        f'<c:symbol val="circle"/>'
        f'<c:size val="8"/>'
        f'<c:spPr>'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:ln w="19050">'
        f'<a:solidFill><a:srgbClr val="{DOT_COLOR}"/></a:solidFill>'
        f'</a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        f'</c:marker>'
        f'<c:xVal>{cat_ref}</c:xVal>'
        f'<c:yVal>{_num_ref(currents)}</c:yVal>'
        f'<c:smooth val="0"/>'
        f'</c:ser>'

        f'<c:dLbls>'
        f'<c:showLegendKey val="0"/><c:showVal val="0"/>'
        f'<c:showCatName val="0"/><c:showSerName val="0"/>'
        f'<c:showPercent val="0"/><c:showBubbleSize val="0"/>'
        f'</c:dLbls>'
        f'<c:axId val="{CAT_AX_ID}"/>'
        f'<c:axId val="{VAL_AX_ID}"/>'
        f'</c:scatterChart>'
    )

    # ── Single median scatter series with uniform X-direction error bars (gold) ─
    valid_medians = [
        m if (m is not None and not (isinstance(m, float) and (math.isnan(m) or math.isinf(m))))
        else None
        for m in medians
    ]
    if any(v is not None for v in valid_medians):
        chart_xml_median = (
            f'<c:scatterChart>'
            f'<c:scatterStyle val="lineMarker"/>'
            f'<c:varyColors val="0"/>'
            f'<c:ser>'
            f'<c:idx val="3"/><c:order val="3"/>'
            f'<c:spPr>'
            f'<a:ln w="25400" cap="rnd"><a:noFill/><a:round/></a:ln>'
            f'<a:effectLst/>'
            f'</c:spPr>'
            f'<c:marker><c:symbol val="none"/></c:marker>'
            f'<c:errBars>'
            f'<c:errDir val="x"/>'
            f'<c:errBarType val="both"/>'
            f'<c:errValType val="fixedVal"/>'
            f'<c:noEndCap val="1"/>'
            f'<c:val val="0.1"/>'
            f'<c:spPr>'
            f'<a:noFill/>'
            f'<a:ln w="28575" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill>'
            f'<a:srgbClr val="{MEDIAN_COLOR}">'
            f'<a:alpha val="80000"/>'
            f'</a:srgbClr>'
            f'</a:solidFill>'
            f'<a:round/>'
            f'</a:ln>'
            f'<a:effectLst/>'
            f'</c:spPr>'
            f'</c:errBars>'
            f'<c:xVal>{cat_ref}</c:xVal>'
            f'<c:yVal>{_num_ref(valid_medians)}</c:yVal>'
            f'<c:smooth val="0"/>'
            f'</c:ser>'
            f'<c:dLbls>'
            f'<c:showLegendKey val="0"/><c:showVal val="0"/>'
            f'<c:showCatName val="0"/><c:showSerName val="0"/>'
            f'<c:showPercent val="0"/><c:showBubbleSize val="0"/>'
            f'</c:dLbls>'
            f'<c:axId val="{CAT_AX_ID}"/>'
            f'<c:axId val="{VAL_AX_ID}"/>'
            f'</c:scatterChart>'
        )
    else:
        chart_xml_median = ''

    # ── Axes + chart closing ─────────────────────────────────────────────────
    chart_xml_tail = (
        # ── Category axis (shared by both chart types) ──────────────────────
        f'<c:catAx>'
        f'<c:axId val="{CAT_AX_ID}"/>'
        f'<c:scaling><c:orientation val="minMax"/></c:scaling>'
        f'<c:delete val="0"/>'
        f'<c:axPos val="b"/>'
        f'<c:numFmt formatCode="General" sourceLinked="1"/>'
        f'<c:majorTickMark val="none"/>'
        f'<c:minorTickMark val="none"/>'
        f'<c:tickLblPos val="low"/>'
        f'<c:spPr>'
        f'<a:noFill/>'
        f'<a:ln w="9525" cap="flat" cmpd="sng" algn="ctr">'
        f'<a:solidFill>'
        f'<a:schemeClr val="tx1">'
        f'<a:lumMod val="15000"/><a:lumOff val="85000"/>'
        f'</a:schemeClr>'
        f'</a:solidFill>'
        f'<a:round/>'
        f'</a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        f'<c:txPr>'
        f'<a:bodyPr rot="-2700000" spcFirstLastPara="1" vertOverflow="ellipsis"'
        f' vert="horz" wrap="square" anchor="ctr" anchorCtr="1"/>'
        f'<a:lstStyle/>'
        f'<a:p><a:pPr><a:defRPr sz="1000" b="0" i="0" u="none" strike="noStrike"'
        f' kern="1200" baseline="0">'
        f'<a:solidFill><a:schemeClr val="tx1">'
        f'<a:lumMod val="65000"/><a:lumOff val="35000"/>'
        f'</a:schemeClr></a:solidFill>'
        f'<a:latin typeface="+mn-lt"/>'
        f'</a:defRPr></a:pPr><a:endParaRPr lang="de-DE"/></a:p>'
        f'</c:txPr>'
        f'<c:crossAx val="{VAL_AX_ID}"/>'
        f'<c:crosses val="autoZero"/>'
        f'<c:auto val="1"/>'
        f'<c:lblAlgn val="ctr"/>'
        f'<c:lblOffset val="100"/>'
        f'<c:noMultiLvlLbl val="0"/>'
        f'</c:catAx>'

        # ── Value axis ──────────────────────────────────────────────────────
        f'<c:valAx>'
        f'<c:axId val="{VAL_AX_ID}"/>'
        f'<c:scaling><c:orientation val="minMax"/></c:scaling>'
        f'<c:delete val="0"/>'
        f'<c:axPos val="l"/>'
        f'<c:majorGridlines>'
        f'<c:spPr>'
        f'<a:ln w="9525" cap="flat" cmpd="sng" algn="ctr">'
        f'<a:solidFill>'
        f'<a:schemeClr val="tx1">'
        f'<a:lumMod val="15000"/><a:lumOff val="85000"/>'
        f'</a:schemeClr>'
        f'</a:solidFill>'
        f'<a:round/>'
        f'</a:ln>'
        f'<a:effectLst/>'
        f'</c:spPr>'
        f'</c:majorGridlines>'
        f'<c:numFmt formatCode="{num_fmt}" sourceLinked="0"/>'
        f'<c:majorTickMark val="none"/>'
        f'<c:minorTickMark val="none"/>'
        f'<c:tickLblPos val="nextTo"/>'
        f'<c:spPr>'
        f'<a:noFill/><a:ln><a:noFill/></a:ln><a:effectLst/>'
        f'</c:spPr>'
        f'<c:txPr>'
        f'<a:bodyPr rot="-60000000" spcFirstLastPara="1" vertOverflow="ellipsis"'
        f' vert="horz" wrap="square" anchor="ctr" anchorCtr="1"/>'
        f'<a:lstStyle/>'
        f'<a:p><a:pPr><a:defRPr sz="1000" b="0" i="0" u="none" strike="noStrike"'
        f' kern="1200" baseline="0">'
        f'<a:solidFill><a:schemeClr val="tx1">'
        f'<a:lumMod val="65000"/><a:lumOff val="35000"/>'
        f'</a:schemeClr></a:solidFill>'
        f'<a:latin typeface="+mn-lt"/>'
        f'</a:defRPr></a:pPr><a:endParaRPr lang="de-DE"/></a:p>'
        f'</c:txPr>'
        f'<c:crossAx val="{CAT_AX_ID}"/>'
        f'<c:crosses val="autoZero"/>'
        f'<c:crossBetween val="between"/>'
        f'</c:valAx>'

        f'</c:plotArea>'
        f'<c:plotVisOnly val="1"/>'
        f'<c:dispBlanksAs val="gap"/>'
        f'<c:showDLblsOverMax val="0"/>'
        f'</c:chart>'
    )
    chart_xml = chart_xml_head + chart_xml_median + chart_xml_tail

    # ── Insert a dummy line chart so python-pptx wires up the chart part ────
    from pptx.chart.data import CategoryChartData as _CCD
    dummy = _CCD()
    dummy.categories = ['x']
    dummy.add_series('s', [1])
    frame = chart_pl.insert_chart(XL_CHART_TYPE.LINE_MARKERS, dummy)
    chart = frame.chart

    # Replace the <c:chart> element with our hand-crafted XML
    chart_space = chart._element
    c_uri = f'{{{C_NS}}}'
    old_chart_el = chart_space.find(f'{c_uri}chart')
    new_chart_el = etree.fromstring(chart_xml)
    if old_chart_el is not None:
        idx = list(chart_space).index(old_chart_el)
        chart_space.remove(old_chart_el)
        chart_space.insert(idx, new_chart_el)
    else:
        chart_space.append(new_chart_el)

    return True


def _add_scatter_chart(traces: list, item: dict, chart_pl) -> bool:
    """
    Insert a line chart into a PowerPoint chart placeholder.
    `traces` is a list of Plotly-like trace dicts with 'x', 'y', 'name'.
    Mirrors the original _add_scatter_chart function.
    """
    scatter_traces = [
        tr for tr in traces
        if tr.get("x") and tr.get("y") and tr.get("showlegend", True)
    ]
    if not scatter_traces:
        return False

    # Build long DataFrame then pivot
    rows = []
    for tr in scatter_traces:
        for x, y in zip(tr["x"], tr["y"]):
            rows.append({"Date": x, "Series": tr["name"], "Value": y})
    df_long = pd.DataFrame(rows)
    if df_long.empty:
        return False

    # sort=False preserves the original x-axis order from the traces.
    # Without this, pandas sorts the "Date" index alphabetically, which
    # breaks non-date x-axes like yield curve maturities (2J, 5J, 10J…
    # would become 10J, 20J, 2J, 30J, 5J in alphabetical order).
    df = df_long.groupby(["Date", "Series"], sort=False)["Value"].first().unstack()
    col_order = [tr["name"] for tr in scatter_traces]
    existing_cols = [c for c in col_order if c in df.columns]
    df = df[existing_cols].reset_index()
    if "Date" in df.columns:
        df = df.set_index("Date")
    # Restore the original x-axis order (unstack may re-sort).
    x_order = [x for x in scatter_traces[0]["x"] if x in df.index] if scatter_traces else []
    if x_order:
        df = df.reindex(x_order)
    try:
        df.index = pd.to_datetime(df.index)
    except Exception:
        pass

    df_interp = df.copy()
    try:
        df_interp = df.interpolate(method="linear")
    except Exception:
        pass

    chart_data = CategoryChartData()
    chart_data.categories = df.index.tolist()
    for col in df.columns:
        chart_data.add_series(col, _clean_chart_values(df_interp[col].tolist()))

    chart_type_enum = XL_CHART_TYPE.LINE_MARKERS
    frame = chart_pl.insert_chart(chart_type_enum, chart_data)
    chart = frame.chart

    for series in chart.series:
        series.format.line.width = Pt(3)
        series.format.line.fill.solid()
        series.marker.style = XL_MARKER_STYLE.NONE

    if len(df.columns) > 1:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
    else:
        chart.has_legend = False

    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.has_title = False
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(130, 135, 150)

    try:
        all_vals = df.values.flatten()
        all_vals = [v for v in all_vals if v is not None and not (isinstance(v, float) and math.isnan(v))]
        if all_vals:
            maxv = float(max(all_vals))
            minv = float(min(all_vals))
            prec = _set_precision(maxv - minv)
            chart.value_axis.minimum_scale = float(_round_to_precision(minv, prec, "down"))
            chart.value_axis.maximum_scale = float(_round_to_precision(maxv, prec, "up"))
    except Exception:
        pass

    sub = item.get("subheading", "")
    y_axis_label = item.get("yAxisLabel", "")
    try:
        all_vals2 = df.values.flatten()
        all_vals2 = [v for v in all_vals2 if v is not None and not (isinstance(v, float) and math.isnan(v))]
        distance = max(abs(min(all_vals2)), abs(max(all_vals2))) if all_vals2 else 100
    except Exception:
        distance = 100
    if y_axis_label == "%" or "Prozent" in sub:
        fmt = "#,##0" if distance > 20 else "#,##0.0"
    else:
        fmt = "#,##0.0" if distance < 20 else "#,##0"
    _set_value_format(chart, fmt)

    # Set y-axis title from metricsConfig yAxisLabel
    _set_value_axis_title(chart, y_axis_label)

    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    if df.index is not None and len(df.index) > 1:
        try:
            days_diff = (df.index[-1] - df.index[0]).days
            if days_diff > 2000:
                _SubElement(chart.category_axis.format.element, "c:majorUnit", val="2")
                _SubElement(chart.category_axis.format.element, "c:majorTimeUnit", val="years")
                _set_date_format(chart, "yyyy")
            elif days_diff > 1000:
                _SubElement(chart.category_axis.format.element, "c:majorTimeUnit", val="years")
                _set_date_format(chart, "yyyy")
            elif days_diff < 60:
                _SubElement(chart.category_axis.format.element, "c:majorUnit", val="5")
                _set_date_format(chart, r"m\/d;@")
            else:
                _SubElement(chart.category_axis.format.element, "c:majorTimeUnit", val="months")
                _set_date_format(chart, r"m\/yy;@" if days_diff > 370 else "mmm.")
        except Exception:
            pass

    _set_cross_between(chart)
    return True


def _get_layouts():
    """Load template and parse layout map (mirrors original get_layouts)."""
    template_path = os.path.abspath(_TEMPLATE_PATH)
    prs = Presentation(template_path)
    layouts = prs.slide_layouts
    layout_map = {}

    for i, lay in enumerate(layouts):
        count = 0
        chart_ph_indices = []
        text_ph_indices = []
        source_ph_indices = []
        text_ph_widths = []
        text_ph_heights = []
        text_ph_idents = []

        for shape in lay.shapes:
            name_parts = shape.name.split(" ")
            kind = name_parts[0]
            if kind == "Diagrammplatzhalter":
                chart_ph_indices.append(shape.placeholder_format.idx)
                count += 1
            elif kind == "Textplatzhalter":
                text_ph_indices.append(shape.placeholder_format.idx)
                text_ph_widths.append(shape.width)
                text_ph_heights.append(shape.height)
                text_ph_idents.append(name_parts[1] if len(name_parts) > 1 else str(i))
            elif kind == "Inhaltsplatzhalter":
                source_ph_indices.append(shape.placeholder_format.idx)

        layout_map[i] = {
            count: {
                "Diagrammplatzhalter": {"Index": chart_ph_indices},
                "Textplatzhalter": {
                    "Index": text_ph_indices,
                    "Width": text_ph_widths,
                    "Height": text_ph_heights,
                    "Ident": text_ph_idents,
                },
                "Inhaltsplatzhalter": {"Index": source_ph_indices},
            }
        }

    return layout_map, layouts, prs


# ─────────────────────────────────────────────────────────────────────────────
# PPTX export  (mirrors original build_pptx)
# ─────────────────────────────────────────────────────────────────────────────

def build_pptx(items: list) -> bytes:
    """
    Build a PowerPoint file from a list of chart items using pptx_template.pptx.

    Item shape:
      {
        "title": str,          # full display title
        "pptx_title": str,     # optional short slide title (falls back to title)
        "subheading": str,     # optional subtitle text
        "source": str,         # optional source text
        "group": int,          # slide group (same group = same slide)
        "chartData": list,
        "regions": list[str],
        "xKey": str,
      }
    """
    layout_map, layouts, prs = _get_layouts()

    # Remove all existing slides from template
    for sldNum in range(len(prs.slides))[::-1]:
        rId = prs.slides._sldIdLst[sldNum].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[sldNum]

    individual_groupings = [it.get("group", 1) for it in items]
    unique_groupings = list(dict.fromkeys(individual_groupings))  # preserve order

    available_sizes = sorted(
        {size for layout_counts in layout_map.values() for size in layout_counts.keys()},
        reverse=True,
    )

    def _process_slide_items(slide_items: list, layout_size: int, layout_index: int):
        """Add one slide for the given items using the given layout."""
        lay = layouts[layout_index]
        slide = prs.slides.add_slide(lay)
        try:
            slide.shapes.title.text = "Add Title"
        except Exception:
            pass

        for c, item in enumerate(slide_items):
            placeholders = slide.placeholders
            try:
                idx_chart  = layout_map[layout_index][layout_size]["Diagrammplatzhalter"]["Index"][c]
                idx_source = layout_map[layout_index][layout_size]["Inhaltsplatzhalter"]["Index"][c]
                text_info  = layout_map[layout_index][layout_size]["Textplatzhalter"]
                text_idxs  = list(text_info["Index"])
                heights    = text_info["Height"]
            except Exception:
                continue

            if layout_size == 1:
                box = heights.index(max(heights))
                del text_idxs[box]
                heading_idx, sub_idx = text_idxs
            elif layout_size == 2:
                box = heights.index(max(heights))
                del text_idxs[box]
                heading_idx = text_idxs[c * 2]
                sub_idx     = text_idxs[c * 2 + 1]
            else:
                heading_idx = text_idxs[2 * c]
                sub_idx     = text_idxs[2 * c + 1]

            placeholders[heading_idx].text = item.get("pptx_title") or item["title"]
            placeholders[sub_idx].text     = item.get("subheading", "")
            placeholders[idx_source].text  = item.get("source", "")

            chart_pl = placeholders[idx_chart]

            # Choose chart renderer based on chart type
            if item.get('chartType') == 'PMIHeatmap':
                # PMI heatmap: replace chart placeholder with a colour-coded table
                _add_pmi_heatmap_table(item, chart_pl, slide)
            elif item.get('chartType') == 'Bar' and item.get('balkenData'):
                # Balken chart: mixed stacked-bar + scatter (range bars + current dot)
                _add_balken_chart(item['balkenData'], item, chart_pl)
            else:
                # Standard time-series line chart
                # Convert our row data to Plotly-like scatter traces
                traces = _chartdata_to_traces(
                    item.get("chartData") or [],
                    item.get("regions") or [],
                    item.get("xKey") or "DatePoint",
                )

                _add_scatter_chart(traces, item, chart_pl)

    for grouping in unique_groupings:
        filtered_items = [it for it in items if it.get("group", 1) == grouping]
        count = len(filtered_items)

        # Find exact layout
        layout_index = next((k for k, v in layout_map.items() if count in v), None)

        if layout_index is not None:
            _process_slide_items(filtered_items, count, layout_index)
        else:
            # Split into chunks that fit available layouts
            items_remaining = list(filtered_items)
            while items_remaining:
                best_size = next((s for s in available_sizes if s <= len(items_remaining)), None)
                if best_size is None:
                    best_size = min(available_sizes)
                slide_batch = items_remaining[:best_size]
                items_remaining = items_remaining[best_size:]
                li = next((k for k, v in layout_map.items() if best_size in v), None)
                if li is None:
                    continue
                _process_slide_items(slide_batch, best_size, li)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
